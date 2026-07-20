#!/usr/bin/env python3
"""Generate the ChessLoop PowerPoint presentation. Run: python3 create_pptx.py"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

IMG_DIR = "/home/nick/dev/lucent/scratchpad/chesslooppix"
GOLD = RGBColor(0xD4, 0xAF, 0x37)
DARK_BG = RGBColor(0x0A, 0x0E, 0x14)
TEXT = RGBColor(0xCB, 0xD5, 0xE1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_bg(slide, color=DARK_BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_title(slide, text, subtitle=None, y=0.5):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GOLD
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
    return box


def add_bullets(slide, items, x=0.6, y=2.0, w=11.5, h=4.5, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▹ {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
    return box


def add_image_slide(slide, img_filename, caption):
    path = os.path.join(IMG_DIR, img_filename)
    assert os.path.exists(path), f"missing image: {path}"
    pic = slide.shapes.add_picture(path, Inches(1.0), Inches(1.7), width=Inches(11.3))
    # cap height so it never overflows the slide
    if pic.height > Inches(5.6):
        ratio = Inches(5.6) / pic.height
        pic.height = Inches(5.6)
        pic.width = int(pic.width * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)


# ---- Slide 1: Title ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
box = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "♞ ChessLoop"
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Self-Hosted Spaced-Repetition Chess Opening Trainer"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0x38, 0xBD, 0xF8)
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Python/FastAPI · React/TypeScript · SQLite · Chessground · Docker"
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT
p3.alignment = PP_ALIGN.CENTER

# ---- Slide 2: Core Loop ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "What ChessLoop Does", "The Core Loop")
add_bullets(s, [
    "Teach — play moves on the Teaching Board; each move saves to a Line in a Library",
    "Practice — the SRS engine (SM-2 variant) picks your weakest positions and quizzes you",
    "Feedback — wrong answers show a red flash + correct-move arrow; you must replay the right move",
    "Track — the Stats page shows accuracy heatmaps by move number and mastery % per opening",
])

# ---- Slide 3: Core Capabilities ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Core Capabilities", "Feature Overview")
add_bullets(s, [
    "Opening Library Builder — teach by playing moves, per-move notes, PGN export",
    "Spaced-Repetition Drill Engine — modified SM-2, leech detection, weakness-biased selection",
    "Public Library Discovery — browse/search/filter, fork, star, read-only viewer",
    "Admin & Security — TOTP MFA, backup/restore, user role management",
])

# ---- Slide 4: Teaching Board ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Teaching Board", "Build Your Repertoire Visually")
add_bullets(s, [
    "Move navigation: ⟪ First · ‹ Previous · Next › · Last ⟫",
    "↓ PGN export with [Event] header",
    "⧭ Duplicate a line (full copy of all moves)",
    "Delete a move and everything after it",
])

# ---- Slide 5: Move Annotations ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Move Annotations", "Document Your Strategy")
add_bullets(s, [
    "Click a move → Add note → type (Ctrl+Enter to save)",
    "Notes persist across every session",
    "Visible in Teaching, Unrated Learning, and Practice modes",
    "Use for: why a move is played, trap warnings, links to analysis",
])

# ---- Slide 6: SRS Algorithm ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Spaced Repetition Engine", "Modified SM-2 with Leech Detection")
add_bullets(s, [
    "Correct: interval grows by ease factor (default 2.5x); Easy +0.15 ease, Hard -0.15",
    "Wrong: interval reset to max(1, interval × 0.25), requeued in 10 minutes",
    "Leech threshold: 4 cumulative wrong answers flags the position",
])

# ---- Slide 7: Practice Sessions ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Practice Sessions", "Drill Your Weaknesses")
add_bullets(s, [
    "Weakest First — auto-starts the most overdue positions",
    "Rated / Unrated — choose whether a session counts toward metrics",
    "Leech Drill — dedicated mode for positions with 4+ wrong answers",
    "Priority: overdue → leeches → new items (20%) → weakness bias (2x under ease 1.8)",
])

# ---- Slide 8: Feedback ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Intelligent Feedback", "Motor Memory Through Physical Execution")
add_bullets(s, [
    "Board flashes red immediately on a wrong answer",
    "Animated arrow shows the correct move",
    "You must physically play it before continuing — recognition alone doesn't count",
])

# ---- Slide 9: Progress Tracking ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Progress Tracking", "See Your Improvement Over Time")
add_bullets(s, [
    "Accuracy heatmap by move number — green ≥80%, yellow 50-79%, red <50%",
    "Mastery badges: not_started → learning → developing → advanced → mastered",
    "Per-library breakdown",
])

# ---- Slide 10: Public Discovery ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Public Library Discovery", "Share and Learn from the Community")
add_bullets(s, [
    "Search/filter by ECO code, color, difficulty",
    "Star favorites, leave comments",
    "Fork / Add to My Library — copies to your account",
    "Learn / View — read-only viewer, no forking required",
])

# ---- Slide 11: Security ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Security & Auth", "JWT + TOTP MFA")
add_bullets(s, [
    "Access token: 15-minute TTL, memory-only storage",
    "Refresh token: 30-day TTL, httpOnly cookie",
    "TOTP MFA via pyotp — Google Authenticator compatible, QR code setup",
])

# ---- Slide 12: Backups ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Backups & Recovery", "Admin Panel, No Downtime")
add_bullets(s, [
    "full — entire database (users, libraries, lines, SRS progress)",
    "content — libraries and lines only, portable to another instance",
    "progress — SRS cards and review log only",
    "10-backup retention, create/download/restore/upload from Admin panel",
])

# ---- Slides 13-17: Screenshots ----
add_image_slide(prs.slides.add_slide(BLANK), "ChessLoop Openings database.png", "Opening Library Browser")
s = prs.slides[-1]
set_bg(s)
add_title(s, "Opening Library Browser", "Your collection at a glance", y=0.3)

s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Learning & Training Arenas", "Read-only study mode", y=0.3)
add_image_slide(s, "ChessLoop Learning and training arenas.png", "Learning Arenas")

s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Tracking Dashboard", "Accuracy heatmaps & mastery badges", y=0.3)
add_image_slide(s, "ChessLoop Tracking dashboard.png", "Tracking Dashboard")

s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Statistics Tracking", "Session-level performance detail", y=0.3)
add_image_slide(s, "Chessloop Statistics tracking.png", "Statistics")

s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Game Tracking & Self-Analysis", "Review every practice session", y=0.3)
add_image_slide(s, "Chessloop Game tracking and self analysis.png", "Game Tracking")

# ---- Slide 18: Appearance ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Appearance & Customization", "Dark/Light Theme, Board & Piece Sets")
add_bullets(s, [
    "Board themes: Brown (default), Blue, Green, Ice, Purple",
    "Piece sets: CBurnett (default), Alpha, Mono, Shadow",
    "App theme: Dark (default) or Light — persists per-user, syncs across devices",
])

# ---- Slide 19: Deployment ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Self-Hosted Deployment", "One Command, Any Linux Box")
add_bullets(s, [
    "curl -fsSL .../deploy.sh | bash",
    "Auto dependency install: git, curl, openssl, Docker",
    "Port scan + conflict-free suggestion",
    "Localhost-only vs. all-interfaces bind choice, health check, systemd autostart",
])

# ---- Slide 20: Seeded Libraries ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "16 Seeded Opening Libraries", "Ready on First Install")
add_bullets(s, [
    "Italian Game — Main Line (White, C50, Beginner)",
    "Ruy López — Closed (White, C84, Intermediate)",
    "Sicilian — Najdorf (Black, B90, Advanced)",
    "Sicilian — Dragon (Black, B70, Advanced)",
    "French Defence — Classical (Black, C11, Intermediate)",
    "+ 11 more, auto-published to Public Discovery",
])

# ---- Slide 21: Tech Stack ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Tech Stack", "Full Breakdown")
add_bullets(s, [
    "Backend — Python 3.12, FastAPI, SQLModel, SQLite (WAL)",
    "Auth — python-jose (JWT), passlib (bcrypt), pyotp (TOTP)",
    "Frontend — React 18, Vite, TypeScript, Chessground, chess.js, TailwindCSS",
    "State — Zustand (auth), TanStack Query (server data)",
    "Infra — Docker Compose, Nginx reverse proxy",
])

# ---- Slide 22: API Surface ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "API Surface", "Full OpenAPI Docs Included")
add_bullets(s, [
    "/api/auth/ — register, login, MFA setup/confirm, refresh",
    "/api/libraries/ — CRUD, publish, fork, active toggle",
    "/api/lines/ — CRUD, moves, per-move notes",
    "/api/practice/ — session start/next/answer/end, due-count",
    "/api/stats/ — heatmap, mastery, leeches, recent sessions",
    "/api/public/ — browse, search, star, comments",
    "/api/admin/ — backups, user management",
], size=16)

# ---- Slide 23: Engineering Highlights ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
add_title(s, "Engineering Highlights", "Notable Design Decisions")
add_bullets(s, [
    "Strict move recording — no fuzzy matching, enforces repertoire discipline",
    "Canonical position keys — normalized FEN enables SRS dedup across transpositions",
    "SQLite WAL mode — concurrent readers + one writer, file = backup",
])

# ---- Slide 24: Closing ----
s = prs.slides.add_slide(BLANK)
set_bg(s)
box = s.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(2.2))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Start Training Today"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Self-hosted, private, built for real spaced-repetition practice."
p2.font.size = Pt(20)
p2.font.color.rgb = TEXT
p2.alignment = PP_ALIGN.CENTER

out_path = "/home/nick/dev/lucent/presentation/chessloop/slides.pptx"
prs.save(out_path)
print(f"Saved {len(prs.slides.__iter__.__self__._sldIdLst)} slides to {out_path}")
